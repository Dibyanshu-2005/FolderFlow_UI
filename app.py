import streamlit as st
import tempfile
import os
from typing import List, Dict
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain.chains import ConversationalRetrievalChain
from langchain.schema import Document
import google.generativeai as genai
from pptx import Presentation
import time
import pyperclip

# Constants
HILTI_RED = "#D00F22"
MAX_QUESTION_LENGTH = 500
FILE_SIZE_LIMIT = 200  # MB

class CustomPPTLoader:
    def __init__(self, file_path: str, original_filename: str):
        self.file_path = file_path
        self.original_filename = original_filename

    def load(self) -> List[Document]:
        try:
            prs = Presentation(self.file_path)
            documents = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
                
                if text_content:
                    text = "\n\n".join(text_content)
                    metadata = {
                        "source": self.original_filename,
                        "slide_number": slide_num,
                        "type": "pptx"
                    }
                    documents.append(Document(page_content=text, metadata=metadata))
            
            return documents
        except Exception as e:
            st.error(f"Error processing PowerPoint file: {str(e)}")
            return []

class DocumentManager:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        self.processed_files = []
        self.qa_chain = None
        self.chat_history = []
        self.total_files = 0
        self.processed_count = 0
        self.start_time = None
        self.file_stats = {
            "pdf": 0,
            "docx": 0,
            "pptx": 0,
            "total_size": 0
        }

    def get_stats(self) -> Dict:
        return {
            "total_files": len(self.processed_files),
            "file_types": self.file_stats,
            "processing_time": time.time() - self.start_time if self.start_time else 0
        }

    def process_file(self, uploaded_file) -> List[Document]:
        try:
            original_filename = uploaded_file.name
            file_size = uploaded_file.size / (1024 * 1024)  # Convert to MB
            self.file_stats["total_size"] += file_size

            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(original_filename)[1]) as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                file_path = tmp_file.name

            documents = []
            if original_filename.endswith('.pdf'):
                loader = PyPDFLoader(file_path)
                documents = loader.load()
                self.file_stats["pdf"] += 1
            elif original_filename.endswith('.docx'):
                loader = Docx2txtLoader(file_path)
                documents = loader.load()
                self.file_stats["docx"] += 1
            elif original_filename.endswith(('.pptx', '.ppt')):
                loader = CustomPPTLoader(file_path, original_filename)
                documents = loader.load()
                self.file_stats["pptx"] += 1
            
            if documents:
                for doc in documents:
                    doc.metadata['source'] = original_filename
                    doc.metadata['size'] = f"{file_size:.2f}MB"

                chunks = self.text_splitter.split_documents(documents)
                self.processed_files.append(original_filename)
                self.processed_count += 1
                return chunks
            return []

        except Exception as e:
            st.error(f"Error processing {original_filename}: {str(e)}")
            return []
        finally:
            if 'file_path' in locals():
                try:
                    os.unlink(file_path)
                except:
                    pass

    def setup_qa_system(self, uploaded_files):
        try:
            self.start_time = time.time()
            all_chunks = []
            self.total_files = len(uploaded_files)
            self.processed_count = 0
            
            progress_bar = st.progress(0)
            status_text = st.empty()
            time_text = st.empty()

            for uploaded_file in uploaded_files:
                file_size = uploaded_file.size / (1024 * 1024)
                if file_size > FILE_SIZE_LIMIT:
                    st.warning(f"File {uploaded_file.name} exceeds size limit of {FILE_SIZE_LIMIT}MB")
                    continue

                status_text.text(f"Processing {uploaded_file.name}...")
                start_time = time.time()
                
                chunks = self.process_file(uploaded_file)
                all_chunks.extend(chunks)
                
                progress = int((self.processed_count / self.total_files) * 100)
                progress_bar.progress(progress)
                
                elapsed_time = time.time() - start_time
                time_text.text(f"Time taken: {elapsed_time:.2f} seconds")
                status_text.text(f"Processed {self.processed_count} of {self.total_files} files ({progress}%)")

            if not all_chunks:
                st.error("No documents were successfully processed!")
                return False

            status_text.text("Setting up QA system...")
            embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
            vectorstore = FAISS.from_documents(all_chunks, embeddings)
            
            self.qa_chain = ConversationalRetrievalChain.from_llm(
                ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=0.7),
                vectorstore.as_retriever(),
                return_source_documents=True
            )
            
            total_time = time.time() - self.start_time
            status_text.text(f"System ready! Total setup time: {total_time:.2f} seconds")
            progress_bar.progress(100)
            return True

        except Exception as e:
            st.error(f"Error setting up QA system: {str(e)}")
            return False

    def ask_question(self, question: str) -> Dict:
        if not self.qa_chain:
            return {"error": "QA system not initialized. Please upload documents first."}
        
        if len(question) > MAX_QUESTION_LENGTH:
            return {"error": f"Question too long. Please limit to {MAX_QUESTION_LENGTH} characters."}
            
        try:
            result = self.qa_chain.invoke({
                "question": question,
                "chat_history": self.chat_history
            })
            
            sources = list({
                doc.metadata.get('source', 'Unknown source')
                for doc in result["source_documents"]
            })
            
            self.chat_history.append((question, result["answer"]))
            return {
                "answer": result["answer"],
                "sources": sources
            }
        except Exception as e:
            return {"error": f"Error processing question: {str(e)}"}

def initialize_session_state():
    if 'manager' not in st.session_state:
        st.session_state.manager = None
    if 'system_ready' not in st.session_state:
        st.session_state.system_ready = False
    if 'messages' not in st.session_state:
        st.session_state.messages = []
    if 'api_key' not in st.session_state:
        st.session_state.api_key = os.getenv('GOOGLE_API_KEY', '')

def copy_to_clipboard(text):
    try:
        pyperclip.copy(text)
        return True
    except:
        return False

def main():
    st.set_page_config(page_title="FolderFlow QA Assistant", page_icon="📚", layout="wide")
    
    # Custom CSS for Hilti branding
    st.markdown(f"""
        <style>
        .stApp {{
            background-color: #f8f9fa;
        }}
        .main-header {{
            color: {HILTI_RED};
            font-weight: bold;
            font-size: 50px;
        }}
        .stButton>button {{
            background-color: {HILTI_RED};
            color: white;
        }}
        </style>
    """, unsafe_allow_html=True)
    
    st.markdown("<h1>FolderFlow for <span class='main-header'>HILTI</span> Technologies</h1>", unsafe_allow_html=True)
    
    initialize_session_state()

    # Two-column layout for main controls
    col1, col2 = st.sidebar.columns(2)
    
    with col1:
        if st.button("Reset Chat", help="Clear chat history"):
            st.session_state.messages = []
            st.rerun()
            
    with col2:
        if st.button("Clear Documents", help="Remove all processed documents"):
            st.session_state.manager = None
            st.session_state.system_ready = False
            st.rerun()

    # API key input in sidebar with persistence
    st.sidebar.header("Configuration")
    api_key = st.sidebar.text_input(
        "Enter Google API Key",
        value=st.session_state.api_key,
        type="password",
        key="api_key_input",
        help="Your Google API key for Generative AI services"
    )
    
    if api_key:
        st.session_state.api_key = api_key
        os.environ["GOOGLE_API_KEY"] = api_key
        genai.configure(api_key=api_key)
    
    # File upload with improved UI
    st.sidebar.header("Upload Documents")
    uploaded_files = st.sidebar.file_uploader(
        "Upload your documents",
        type=['pdf', 'docx', 'pptx', 'ppt'],
        accept_multiple_files=True,
        help=f"Upload PDF, Word, or PowerPoint files (max {FILE_SIZE_LIMIT}MB each)"
    )

    if uploaded_files and api_key and st.sidebar.button("Process Documents", type="primary"):
        with st.spinner("Processing documents..."):
            manager = DocumentManager()
            if manager.setup_qa_system(uploaded_files):
                st.session_state.manager = manager
                st.session_state.system_ready = True
                st.sidebar.success("✅ Documents processed successfully!")
            else:
                st.sidebar.error("❌ Setup failed. Please try again.")

    # Display document statistics in sidebar
    if st.session_state.system_ready and st.session_state.manager:
        stats = st.session_state.manager.get_stats()
        st.sidebar.header("Document Statistics")
        st.sidebar.metric("Total Files", stats["total_files"])
        st.sidebar.metric("Total Size", f"{stats['file_types']['total_size']:.2f}MB")
        
        # File type breakdown
        st.sidebar.header("File Types")
        for file_type, count in stats["file_types"].items():
            if file_type != "total_size":
                st.sidebar.text(f"{file_type.upper()}: {count}")

    # Chat interface with enhancements
    if st.session_state.system_ready:
        # Display chat history with copy buttons
        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.write(message["content"])
                if "sources" in message:
                    st.markdown("**Sources:**")
                    for source in message["sources"]:
                        st.markdown(f"- {source}")
                
                # Add copy button for responses
                if message["role"] == "assistant":
                    if st.button("📋 Copy", key=f"copy_{len(st.session_state.messages)}"):
                        if copy_to_clipboard(message["content"]):
                            st.success("Copied to clipboard!")
                        else:
                            st.error("Failed to copy. Please try manually.")

        # Chat input with character counter
        prompt = st.chat_input("Ask a question about your documents")
        if prompt:
            char_count = len(prompt)
            if char_count > MAX_QUESTION_LENGTH:
                st.warning(f"Question too long! ({char_count}/{MAX_QUESTION_LENGTH} characters)")
            else:
                st.session_state.messages.append({"role": "user", "content": prompt})
                with st.chat_message("user"):
                    st.write(prompt)

                with st.chat_message("assistant"):
                    with st.spinner("Thinking..."):
                        response = st.session_state.manager.ask_question(prompt)
                        if "error" in response:
                            st.error(response["error"])
                        else:
                            st.write(response["answer"])
                            if response["sources"]:
                                st.markdown("**Sources:**")
                                for source in response["sources"]:
                                    st.markdown(f"- {source}")
                            
                            st.session_state.messages.append({
                                "role": "assistant",
                                "content": response["answer"],
                                "sources": response["sources"]
                            })
    else:
        st.info("👋 Hi, Welcome to our prototype FolderFlow - More than just 'search' \n\n"
                "Prepared by Dibyanshu and Sajjad - Happy Exploring!\n\n"
                "Kindly follow these steps to begin:\n"
                "1. Enter your Google API key in the sidebar\n"
                "2. Upload your documents\n"
                "3. Click 'Process Documents' to start")

if __name__ == "__main__":
    main()
